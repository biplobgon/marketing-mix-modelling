"""
src/reports/generate_pptx.py
=============================
Automated generation of a BCG-style 15-slide MMM stakeholder deck.

Uses python-pptx to programmatically build a professionally structured
PowerPoint presentation suitable for C-suite and board-level audiences.

Usage:
    python src/reports/generate_pptx.py
    # OR from notebooks:
    from src.reports.generate_pptx import generate_deck
    generate_deck(config, outputs_dir="outputs/")

Output:
    outputs/reports/MMM_Stakeholder_Deck.pptx

Slide structure (15 slides):
    1  Cover — Project Title & Tagline
    2  Executive Summary — 3-box "So What"
    3  Business Problem & Data Landscape
    4  MMM Methodology Overview & DAG
    5  Data & Channels in Scope
    6  Adstock & Saturation Effects
    7  Framework Comparison — 4-model ROAS side-by-side
    8  Sales Decomposition Waterfall
    9  Channel Contribution & ROAS Ranked
    10 Marginal ROAS & Saturation Frontier
    11 Budget Optimisation Results
    12 Scenario Comparison Table
    13 Strategic Recommendations (2×2 Efficiency vs Scale Matrix)
    14 Risks & Model Limitations
    15 Appendix — Technical Notes & Model Equations
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------------
# Colour palette (BCG-inspired)
# ---------------------------------------------------------------------------
PALETTE = {
    "dark_green":  RGBColor(0x00, 0x6A, 0x4E),   # BCG primary
    "light_green": RGBColor(0x72, 0xBF, 0x44),
    "dark_grey":   RGBColor(0x3D, 0x3D, 0x3D),
    "mid_grey":    RGBColor(0x80, 0x80, 0x80),
    "light_grey":  RGBColor(0xF2, 0xF2, 0xF2),
    "white":       RGBColor(0xFF, 0xFF, 0xFF),
    "accent_blue": RGBColor(0x00, 0x5B, 0x96),
}

# Slide dimensions: widescreen 16:9
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helper: add title text box
# ---------------------------------------------------------------------------

def _add_title(slide, text: str, top: float = 0.4, font_size: int = 28) -> None:
    """Add a styled slide title text box.

    Args:
        slide: python-pptx slide object.
        text: Title text.
        top: Distance from top in inches.
        font_size: Font size in points.

    TODO:
        - Add text box: left=0.5", top=top, width=12", height=0.8"
        - Set font: Bold, font_size, PALETTE["dark_green"]
        - Set alignment: PP_ALIGN.LEFT
    """
    # TODO: implement _add_title
    pass


# ---------------------------------------------------------------------------
# Helper: add body text box
# ---------------------------------------------------------------------------

def _add_body(
    slide,
    text: str,
    left: float = 0.5,
    top: float = 1.5,
    width: float = 12.0,
    height: float = 5.0,
    font_size: int = 14,
) -> None:
    """Add a body text box with given dimensions and content.

    Args:
        slide: python-pptx slide object.
        text: Body text (use \\n for newlines).
        left: Left margin in inches.
        top: Top margin in inches.
        width: Width in inches.
        height: Height in inches.
        font_size: Font size in points.

    TODO:
        - Add text box at (left, top, width, height)
        - Set word wrap True
        - Set font size and colour PALETTE["dark_grey"]
    """
    # TODO: implement _add_body
    pass


# ---------------------------------------------------------------------------
# Helper: add image placeholder
# ---------------------------------------------------------------------------

def _add_image_placeholder(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    label: str = "Chart",
) -> None:
    """Add a grey placeholder box for charts not yet generated.

    Args:
        slide: python-pptx slide object.
        left, top, width, height: Position and size in inches.
        label: Text label displayed inside the placeholder.

    TODO:
        - Add rectangle shape with fill PALETTE["light_grey"]
        - Add centred label text in PALETTE["mid_grey"]
    """
    # TODO: implement image placeholder
    pass


# ---------------------------------------------------------------------------
# Slide builders (one function per slide)
# ---------------------------------------------------------------------------

def _slide_01_cover(prs: Presentation) -> None:
    """Slide 1 — Cover: Project Title & Tagline.

    TODO:
        - Add full-bleed background rectangle in PALETTE["dark_green"]
        - Add title: "Marketing Mix Modelling" (white, 40pt)
        - Add subtitle: tagline from README (white, 18pt)
        - Add date: "April 2026" (light_grey, 14pt)
        - Add BCG-style decorative bar at bottom
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    # TODO: implement slide 1
    _add_title(slide, "Marketing Mix Modelling", top=2.5, font_size=40)


def _slide_02_exec_summary(prs: Presentation, findings: dict[str, str]) -> None:
    """Slide 2 — Executive Summary: 3-box "So What".

    Args:
        findings: Dict with keys "insight_1", "insight_2", "insight_3".

    TODO:
        - Title: "Executive Summary"
        - Three equal-width boxes with BCG "So What" framing:
          Box 1: Media accounts for 35–45% of revenue
          Box 2: Digital channels deliver highest ROAS
          Box 3: 22–25% uplift available through reallocation
        - Each box: dark_green header + body text
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Executive Summary")
    # TODO: implement 3-box layout


def _slide_03_business_problem(prs: Presentation) -> None:
    """Slide 3 — Business Problem & Data Landscape.

    TODO:
        - Problem statement text (left 60%)
        - Data landscape summary box (right 40%):
          156 weeks, 8 geos, 6 channels, 3 brands
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "The Business Problem & Data Landscape")
    # TODO: implement slide 3


def _slide_04_methodology(prs: Presentation) -> None:
    """Slide 4 — MMM Methodology Overview & DAG.

    TODO:
        - Left panel: MMM formula (rendered as text/math notation)
        - Right panel: DAG diagram (image or placeholder)
        - Brief text: Adstock + Saturation explanation
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "MMM Methodology Overview")
    # TODO: implement slide 4


def _slide_05_data_channels(prs: Presentation) -> None:
    """Slide 5 — Data & Channels in Scope.

    TODO:
        - Channel list table: TV, YouTube, Facebook, Instagram, Print, Radio
        - Geo coverage map placeholder
        - Date range and row count summary
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Data & Channels in Scope")
    # TODO: implement slide 5


def _slide_06_adstock_saturation(prs: Presentation) -> None:
    """Slide 6 — Adstock & Saturation Effects.

    TODO:
        - Left: adstock decay curve images (outputs/figures/adstock_*.png)
        - Right: saturation curves (outputs/figures/saturation_*.png)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Adstock & Saturation Effects")
    _add_image_placeholder(slide, 0.5, 1.5, 6.0, 5.0, "Adstock Decay Curves")
    _add_image_placeholder(slide, 6.8, 1.5, 6.0, 5.0, "Saturation Curves")
    # TODO: replace placeholders with actual figure images


def _slide_07_framework_comparison(prs: Presentation) -> None:
    """Slide 7 — Framework Comparison: 4-model ROAS side-by-side.

    TODO:
        - Insert ROAS comparison bar chart (outputs/figures/roas_comparison.png)
        - Per-framework summary text boxes
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Framework Comparison — Channel ROAS")
    _add_image_placeholder(slide, 0.5, 1.5, 12.0, 5.0, "ROAS Comparison Chart")
    # TODO: replace with actual figure


def _slide_08_decomposition(prs: Presentation) -> None:
    """Slide 8 — Sales Decomposition Waterfall.

    TODO:
        - Insert waterfall chart (outputs/figures/decomposition_waterfall.png)
        - Key stat: "XX% of sales are media-driven"
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Sales Decomposition Waterfall")
    _add_image_placeholder(slide, 0.5, 1.5, 12.0, 5.0, "Sales Decomposition Waterfall")
    # TODO: replace with actual figure


def _slide_09_roas_ranked(prs: Presentation) -> None:
    """Slide 9 — Channel Contribution & ROAS Ranked.

    TODO:
        - ROAS ranked bar chart (horizontal)
        - Call-out box: highest ROAS channel
        - Call-out box: most saturated channel
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Channel Contribution & ROAS Ranked")
    _add_image_placeholder(slide, 0.5, 1.5, 12.0, 5.0, "Channel ROAS Bar Chart")


def _slide_10_mroas_saturation(prs: Presentation) -> None:
    """Slide 10 — Marginal ROAS & Saturation Frontier.

    TODO:
        - mROAS curve grid per channel (small multiples)
        - Saturation frontier scatter: ROAS vs. spend
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Marginal ROAS & Saturation Frontier")
    _add_image_placeholder(slide, 0.5, 1.5, 6.0, 5.0, "mROAS Curves")
    _add_image_placeholder(slide, 6.8, 1.5, 6.0, 5.0, "Saturation Frontier")


def _slide_11_budget_optimisation(prs: Presentation) -> None:
    """Slide 11 — Budget Optimisation Results.

    TODO:
        - Current vs. optimal allocation bar chart
        - Revenue uplift KPI box: "+22–25% incremental revenue"
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Budget Optimisation Results")
    _add_image_placeholder(slide, 0.5, 1.5, 9.0, 5.0, "Current vs. Optimal Allocation")


def _slide_12_scenario_comparison(prs: Presentation) -> None:
    """Slide 12 — Scenario Comparison Table.

    TODO:
        - Table: Scenario | Total Budget | TV% | Digital% | Projected Revenue | Uplift
        - Three rows: Flat / +10% / +25%
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Scenario Comparison")
    # TODO: add pptx table with scenario data


def _slide_13_recommendations(prs: Presentation) -> None:
    """Slide 13 — Strategic Recommendations (2×2 Matrix).

    TODO:
        - 2×2 matrix: Efficiency (ROAS) on Y-axis, Scale (spend) on X-axis
        - Quadrants: Underinvested / Scale Up / Harvest / Reduce
        - Channel bubbles positioned by ROAS and spend
        - 3 strategic bullets below the matrix
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Strategic Recommendations")
    _add_image_placeholder(slide, 0.5, 1.5, 8.0, 5.0, "2×2 Efficiency vs Scale Matrix")


def _slide_14_risks(prs: Presentation) -> None:
    """Slide 14 — Risks & Model Limitations.

    TODO:
        - Bulleted risk list: collinearity, data quality, external validity
        - Mitigation column for each risk
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Risks & Model Limitations")
    risks = [
        "Collinearity between channels — mitigated by Bayesian priors & Ridge",
        "Synthetic data — directional findings only, not field-calibrated",
        "No incrementality / lift-test validation yet",
        "Meridian & PyMC uncertainty ranges ±15–20% on ROAS estimates",
        "Creative quality proxied by content score (AI-generated, not A/B tested)",
    ]
    _add_body(slide, "\n".join(f"• {r}" for r in risks))


def _slide_15_appendix(prs: Presentation) -> None:
    """Slide 15 — Appendix: Technical Notes & Model Equations.

    TODO:
        - MMM equation typeset
        - Adstock formula and decay-to-half-life table
        - Hill saturation formula
        - Framework hyperparameter summary
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Appendix — Technical Notes")
    _add_body(
        slide,
        "Sales_t = α + Σ β_m · S(Adstock(x_m,t)) + Σ γ_c · z_c,t + ε_t\n\n"
        "Adstock: x*_t = x_t + α · x*_{t-1}\n"
        "Hill:    S(x) = x^β / (x^β + K^β)\n\n"
        "See notebooks 03 and 04 for full parameter derivations.",
        font_size=12,
    )


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def generate_deck(
    config: dict[str, Any] | None = None,
    outputs_dir: str = "outputs/",
    findings: dict[str, str] | None = None,
) -> str:
    """Generate the full 15-slide BCG-style stakeholder deck.

    Args:
        config: Configuration dictionary. Loaded from default path if None.
        outputs_dir: Root directory for outputs (figures + reports).
        findings: Optional dict of key findings to populate slide 2.
                  Falls back to README defaults if None.

    Returns:
        Path to the generated .pptx file.

    TODO:
        - Load config if not provided
        - Load default findings if not provided
        - Instantiate pptx.Presentation()
        - Set slide dimensions to 16:9 widescreen
        - Call all 15 _slide_XX_*() builder functions in order
        - Save to outputs_dir/reports/MMM_Stakeholder_Deck.pptx
        - Print confirmation message
        - Return output path
    """
    # TODO: load config
    if config is None:
        from src.data_prep import load_config
        config = load_config()

    if findings is None:
        findings = {
            "insight_1": "35–45% of sales are media-driven; 55–65% is structural baseline.",
            "insight_2": "Facebook + Instagram deliver the highest ROAS; TV is approaching saturation.",
            "insight_3": "Reallocating ~15% of TV budget to digital delivers +22–25% revenue uplift.",
        }

    # TODO: create and populate presentation
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # TODO: call slide builders in order
    _slide_01_cover(prs)
    _slide_02_exec_summary(prs, findings)
    _slide_03_business_problem(prs)
    _slide_04_methodology(prs)
    _slide_05_data_channels(prs)
    _slide_06_adstock_saturation(prs)
    _slide_07_framework_comparison(prs)
    _slide_08_decomposition(prs)
    _slide_09_roas_ranked(prs)
    _slide_10_mroas_saturation(prs)
    _slide_11_budget_optimisation(prs)
    _slide_12_scenario_comparison(prs)
    _slide_13_recommendations(prs)
    _slide_14_risks(prs)
    _slide_15_appendix(prs)

    # TODO: save file
    output_path = Path(outputs_dir) / "reports" / "MMM_Stakeholder_Deck.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"✅ Deck saved to: {output_path}")
    return str(output_path)


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    generate_deck()
